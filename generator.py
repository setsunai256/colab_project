
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PPTX_RGBColor
from docx.shared import RGBColor
from docx.oxml.ns import qn
import os
import traceback
from transformers import CLIPTokenizer

# Токенизатор CLIP для ограничения длины prompt
clip_tokenizer = CLIPTokenizer.from_pretrained("openai/clip-vit-base-patch32")

def truncate_prompt(prompt: str, max_tokens: int = 77):
    tokens = clip_tokenizer(prompt)["input_ids"]
    if len(tokens) > max_tokens:
        truncated = clip_tokenizer.decode(tokens[:max_tokens], skip_special_tokens=True)
        return truncated
    return prompt

def generate_report(topic, model, tokenizer, generation_config, pipe):
    if not topic:
        return None, None
    try:
        os.makedirs("docs/images", exist_ok=True)

        sections = ['Введение', 'Заключение']
        question = f'Придумай от 2 до 3 кратких, четких и конкретных названий основных разделов для доклада на тему "{topic}". Ответ должен содержать только названия разделов, каждое с новой строки.'
        response, _ = model.chat(tokenizer, None, question, generation_config, history=None, return_history=True)
        additional_sections = [s.strip() for s in response.split('\n') if s.strip()]
        sections[1:1] = additional_sections

        doc = Document()
        style = doc.styles['Normal']
        style.font.name = 'Times New Roman'
        style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Times New Roman')
        style.font.size = Pt(14)
        style.paragraph_format.first_line_indent = Cm(1.25)
        style.paragraph_format.line_spacing = 1.5
        style.paragraph_format.space_after = Pt(0)

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic

        for section in sections:
            question = f'Напиши развернутый текст для раздела "{section}" доклада на тему "{topic}".'
            response, _ = model.chat(tokenizer, None, question, generation_config, history=None, return_history=True)
            response = response.strip() or f'Раздел "{section}" временно не заполнен.'

            heading = doc.add_heading(section, level=1)
            run = heading.runs[0]
            run.font.name = 'Times New Roman'
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0, 0, 0)
            heading.alignment = PP_ALIGN.CENTER

            paragraph = doc.add_paragraph(response)
            paragraph.alignment = PP_ALIGN.JUSTIFY

            slide_question = f'Сформулируй тезисный текст для слайда раздела "{section}" доклада на тему "{topic}". Не более 5 коротких пунктов.'
            slide_response, _ = model.chat(tokenizer, None, slide_question, generation_config, history=None, return_history=True)
            slide_response = slide_response.strip()

            prompt_question = f'Generate a detailed and visually descriptive prompt in English for an image related to the section "{section}" of the report on the topic "{topic}".'
            image_prompt, _ = model.chat(tokenizer, None, prompt_question, generation_config, history=None, return_history=True)
            image_prompt = truncate_prompt(image_prompt)

            image = pipe(image_prompt).images[0]
            image_path = f"docs/images/{section.replace(' ', '_')}.png"
            image.save(image_path)

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = section

            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
            text_frame = textbox.text_frame
            text_frame.text = slide_response if slide_response else 'Тезисы отсутствуют.'

            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4), height=Inches(4.5))

        final_slide = prs.slides.add_slide(prs.slide_layouts[1])
        final_slide.shapes.title.text = "Спасибо за внимание!"
        final_slide.shapes.placeholders[1].text = "Буду рад ответить на ваши вопросы."

        doc_name = f"docs/Доклад_{topic.replace(' ', '_')}.docx"
        ppt_name = f"docs/Презентация_{topic.replace(' ', '_')}.pptx"
        doc.save(doc_name)
        prs.save(ppt_name)

        return doc_name, ppt_name

    except Exception as e:
        print("Ошибка в generate_report:", e)
        print(traceback.format_exc())
        return None, None

